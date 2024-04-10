import logging
from abc import ABC, abstractmethod

import cloudscraper
from bs4 import BeautifulSoup
from langchain_anthropic import ChatAnthropic
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi

from method_type import MethodType

logger = logging.getLogger(__name__)

PROMPT_TEXT_SUMMARIER = """
以下の文章の要点を抽出して、それぞれに対して詳細に説明をして。

制約条件:
- 冒頭にタイトルを記載する
- 要約以外の情報は不要
- bullet に数字は使わない
- 可能な限り日本語で記述する

品質を上げるヒント:
- 本文の中に具体例がある場合はそれを含めると良い。必要に応じてそのまま引用する。
---
{input}
"""


class HTTPClient:
    def __init__(self) -> None:
        self.client = cloudscraper.create_scraper()

    def get(self, url: str) -> str:
        html_content = self.client.get(url).text
        body_text = BeautifulSoup(html_content, "html.parser").get_text()
        return body_text


class TextSummarizer:
    def __init__(self):
        self.prompt = ChatPromptTemplate.from_template(PROMPT_TEXT_SUMMARIER)
        # max_tokens_to_sample は、default の 1024 だと文章が切れることがあるみたいなので 2000 に設定する。
        # 2000 以下にしないと Discord のメッセージ上限に引っかかる。
        self.model = ChatAnthropic(
            model="claude-3-opus-20240229", temperature=0, max_tokens_to_sample=4096
        )
        self.output_parser = StrOutputParser()
        self.chain = self.prompt | self.model | self.output_parser

    def summarize(self, input):
        return self.chain.invoke({"input": input})


class BaseSummarizer(ABC):
    @abstractmethod
    def summarize(self, url: str) -> str:
        pass


class WebSummarizer(BaseSummarizer):
    def __init__(
        self, text_summarizer: TextSummarizer, http_client: HTTPClient
    ) -> None:
        self.text_summrizer = text_summarizer
        self.http_client = http_client

    def summarize(self, url: str) -> str:
        body_text = self.http_client.get(url)
        return self.text_summrizer.summarize(body_text)


class YouTubeSummarizer(BaseSummarizer):
    def __init__(self, text_summarizer: TextSummarizer) -> None:
        self.text_summrizer = text_summarizer

    def _get_video_id(self, url: str) -> str:
        # いくつか異なる形式の URL に対応する
        # まず https://www.youtube.com/watch?v=xxxxx 形式の URL から video_id を取得する
        if "v=" in url:
            return url.split("v=")[1]
        # 次に https://youtu.be/xxxxx 形式の URL から video_id を取得する
        if "youtu.be" in url:
            return url.split("/")[-1]
        # どちらの形式でもない場合は例外を発生させる
        logger.error(f"Invalid URL format: {url}")
        raise ValueError(f"Invalid URL format: {url}")

    def _get_youtube_content(self, url: str) -> str:
        video_id = self._get_video_id(url)
        transcript = YouTubeTranscriptApi.get_transcript(
            video_id, languages=["ja", "en"]
        )
        content = ""
        for i in transcript:
            content += i["text"]
        return content

    def summarize(self, url: str) -> str:
        content = self._get_youtube_content(url)
        return self.text_summrizer.summarize(content)


class ArXivSummarizer(BaseSummarizer):
    def __init__(
        self, text_summarizer: TextSummarizer, http_client: HTTPClient
    ) -> None:
        self.text_summrizer = text_summarizer
        self.http_client = http_client

    def _modify_arxiv_url(self, url: str) -> str:
        """
        https://ar5iv.labs.arxiv.org/
        に習って、 arXiv の URL を変換して HTML のページを見られるようにする

           https://arxiv.org/abs/1910.06709
        -> https://ar5iv.org/abs/1910.06709
        """
        return url.replace("arxiv.org", "ar5iv.org")

    def summarize(self, url: str) -> str:
        modified_url = self._modify_arxiv_url(url)
        body_text = self.http_client.get(modified_url)
        return self.text_summrizer.summarize(body_text)


class PPTSummarizer(BaseSummarizer):
    def __init__(self, text_summarizer: TextSummarizer) -> None:
        self.text_summrizer = text_summarizer

    def extract_text_from_slide(self, slide):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        return text

    def extract_text_from_presentation(self, pptx_path: str):
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            text += self.extract_text_from_slide(slide)
        return text

    def summarize(self, pptx_path: str) -> str:
        text = self.extract_text_from_presentation(pptx_path)
        result = self.text_summrizer.summarize(text)
        return result


class SummarizerBuilder:
    def build_summarizer(self, method: str) -> BaseSummarizer | None:
        text_summarizer = TextSummarizer()
        http_client = HTTPClient()
        summerizer_map = {
            MethodType.WEB.value: WebSummarizer(text_summarizer, http_client),
            MethodType.YOUTUBE.value: YouTubeSummarizer(text_summarizer),
            MethodType.ARXIV.value: ArXivSummarizer(text_summarizer, http_client),
            MethodType.NONE.value: None,
        }
        return summerizer_map[method]
